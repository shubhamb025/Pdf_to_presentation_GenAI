import os
import re
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Define themes with specific files
THEMES = {
    'theme1': {
        'name': 'Theme 1',
        'file': 'Presentation1.pptx'
    },
    'theme2': {
        'name': 'Theme 2',
        'file': 'Presentation2.pptx'
    },
    'theme3': {
        'name': 'Theme 3',
        'file': 'Presentation3.pptx'
    },
    'theme4': {
        'name': 'Theme 4',
        'file': 'Presentation4.pptx'
    }
}

def list_available_themes():
    """Return a list of available themes."""
    return [{'key': key, 'name': theme['name']} for key, theme in THEMES.items()]

def parse_vba_file(vba_file_path):
    slides_data = []
    current_slide = None
    content_buffer = []  # Changed to list to store each point separately

    with open(vba_file_path, 'r', encoding='utf-8') as file:
        for line_number, line in enumerate(file, 1):
            line = line.strip()
            try:
                if "Set sld = ppt.Slides.Add" in line:
                    if current_slide:
                        # Store content as list of points
                        current_slide['content'] = [point for point in content_buffer if point.strip()]
                        slides_data.append(current_slide)
                    current_slide = {'title': '', 'content': []}
                    content_buffer = []  # Reset buffer for new slide
                elif ".Shapes.Title.TextFrame.TextRange.Text =" in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        current_slide['title'] = match.group(1)
                elif ".Text =" in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        # Add each point as a separate item
                        point = match.group(1).strip()
                        if point:
                            content_buffer.append(point)
                elif '& _' in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        # Add continuation text to the last point
                        if content_buffer:
                            content_buffer[-1] += match.group(1)
                        else:
                            content_buffer.append(match.group(1))
            except ValueError as e:
                print(f"Error processing line {line_number}: {line}")
                print(f"Error details: {str(e)}")

    if current_slide:
        # Store content as list of points for the last slide
        current_slide['content'] = [point for point in content_buffer if point.strip()]
        slides_data.append(current_slide)

    return slides_data

def detect_format_type(slides_data):
    """Detect the format type based on the content structure"""
    # Look for indicators in the content
    format_indicators = {
        'qa_format': 0,
        'table_format': 0,
        'diagram_format': 0,
        'bullet_only': 0
    }
    
    # Check slide titles and content for format indicators
    for slide in slides_data:
        # Check for Q&A format
        if slide['title'].endswith('?'):
            format_indicators['qa_format'] += 1
            
        # Get content as list of points
        content_points = slide['content']
        if not isinstance(content_points, list):
            # Convert to list if it's a string
            content_points = [content_points] if content_points else []
            
        # Check each bullet point for format indicators
        for point in content_points:
            point_lower = point.lower() if isinstance(point, str) else ''
            
            # Check for table descriptions
            if 'table' in point_lower and ('row' in point_lower or 'column' in point_lower):
                format_indicators['table_format'] += 1
                
            # Check for diagram descriptions
            if any(term in point_lower for term in ['diagram', 'chart', 'graph', 'visualize']):
                format_indicators['diagram_format'] += 1
                
        # Check if content is short bullet points only
        if content_points and all(isinstance(point, str) and len(point) < 120 for point in content_points):
            format_indicators['bullet_only'] += 1
            
    # Determine the dominant format type
    if format_indicators['qa_format'] > len(slides_data) / 3:
        return 'qa'
    elif format_indicators['table_format'] > len(slides_data) / 3:
        return 'tables'
    elif format_indicators['diagram_format'] > len(slides_data) / 3:
        return 'diagrams'
    elif format_indicators['bullet_only'] > len(slides_data) / 2:
        return 'bullets_only'
    else:
        return 'standard'

def apply_theme_to_slide(slide, theme):
    """Apply theme colors and styles to a slide."""
    # Apply theme to title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.color.rgb = theme['title_color']
        title_frame.paragraphs[0].font.size = theme['font_size']['title']
    
    # Apply theme to body
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if shape != slide.shapes.title:  # Skip title
                    paragraph.font.color.rgb = theme['body_color']
                    paragraph.font.size = theme['font_size']['body']

def get_image_titles():
    """Read image titles from the file."""
    titles_file = os.path.join('extract', 'image_titles.txt')
    titles = {}
    if os.path.exists(titles_file):
        with open(titles_file, 'r', encoding='utf-8') as f:
            for line in f:
                key, title = line.strip().split('|', 1)
                titles[key] = title
    return titles

def create_image_slide(prs, images_dir, page_number):
    """Create a new slide specifically for images."""
    # Check if images directory exists
    if not os.path.exists(images_dir):
        print(f"Images directory not found: {images_dir}")
        return None
        
    # Find all images for this page - check both PDF extraction format and web image format
    slide_images = [f for f in os.listdir(images_dir) 
                   if (f'page_{page_number}_img_' in f.lower() or 
                       f'topic_{page_number}_img_' in f.lower() or
                       f'web_img_{page_number}_' in f.lower()) and 
                   any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
    
    if not slide_images:
        print(f"No images found for page {page_number}")
        return None

    print(f"Found {len(slide_images)} images for page {page_number}")
    
    # Get image titles from the titles file
    image_titles = get_image_titles()
    
    # Sort images by their numbers
    def get_img_number(filename):
        if 'img_' in filename:
            match = re.search(r'img_(\d+)', filename)
            if match:
                return int(match.group(1))
        return 0
    
    slide_images.sort(key=get_img_number)
    
    # Create a new slide for images
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title to the slide
    title_text = f"Visual Content for Section {page_number}"
    
    # Try to get a better title from image titles
    for img_file in slide_images:
        img_key = None
        if 'page_' in img_file:
            img_index = re.search(r'img_(\d+)', img_file).group(1)
            img_key = f"page_{page_number}_img_{img_index}"
        elif 'topic_' in img_file:
            img_index = re.search(r'img_(\d+)', img_file).group(1)
            img_key = f"page_{page_number}_img_{img_index}"
        elif 'web_img_' in img_file:
            img_index = re.search(r'web_img_\d+_(\d+)', img_file).group(1)
            img_key = f"page_{page_number}_img_{img_index}"
            
        if img_key and img_key in image_titles:
            title_text = image_titles[img_key]
            break
    
    # Add the title to the slide
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Calculate layout for images
    if len(slide_images) == 1:
        left_margin = Inches(1)
        top_margin = Inches(1.5)
        img_width = Inches(8)
        img_height = Inches(5.5)
    else:
        left_margin = Inches(0.5)
        top_margin = Inches(1.5)
        img_width = Inches(4.5)
        img_height = Inches(3.5)
        
    # Add images with captions
    for idx, image_file in enumerate(slide_images):
        image_path = os.path.join(images_dir, image_file)
        try:
            print(f"Processing image: {image_file}")
            
            # Calculate position
            row = idx // 2
            col = idx % 2
            left = left_margin + (col * (img_width + Inches(0.5)))
            top = top_margin + (row * (img_height + Inches(0.7)))
            
            # Verify image exists
            if not os.path.exists(image_path):
                print(f"Image file not found: {image_path}")
                continue
                
            # Add the image
            pic = slide.shapes.add_picture(
                image_path,
                left=left,
                top=top,
                width=img_width,
                height=img_height
            )
            print(f"Added image to slide: {image_file}")
            
            # Get image title
            img_key = None
            if 'page_' in image_file:
                img_index = re.search(r'img_(\d+)', image_file).group(1)
                img_key = f"page_{page_number}_img_{img_index}"
            elif 'topic_' in image_file:
                img_index = re.search(r'img_(\d+)', image_file).group(1)
                img_key = f"page_{page_number}_img_{img_index}"
            elif 'web_img_' in image_file:
                img_index = re.search(r'web_img_\d+_(\d+)', image_file).group(1)
                img_key = f"page_{page_number}_img_{img_index}"
            
            # Set caption text
            caption_text = f"Figure {idx+1}"
            if img_key and img_key in image_titles:
                caption_text = image_titles[img_key]
            
            # Add caption
            caption = slide.shapes.add_textbox(
                left=left,
                top=top + img_height + Inches(0.1),
                width=img_width,
                height=Inches(0.5)
            )
            caption_frame = caption.text_frame
            caption_frame.word_wrap = True
            caption_frame.text = caption_text
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            caption_frame.paragraphs[0].font.size = Pt(10)
            caption_frame.paragraphs[0].font.italic = True
            
            print(f"Added caption: {caption_text}")
            
        except Exception as e:
            print(f"Error adding image {image_file}: {str(e)}")
            traceback.print_exc()
    
    return slide

def create_bullet_slide(prs, slide_title, bullet_points, format_type='standard'):
    """Create a slide with properly formatted bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title = slide.shapes.title
    title.text = slide_title
    
    # Get content placeholder
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Format bullet points based on format type
    for i, point in enumerate(bullet_points):
        point_text = point.strip()
        if point_text.startswith('- '):
            point_text = point_text[2:]
            
        p = tf.add_paragraph()
        p.text = point_text
        p.level = 0
        
        # Set line spacing and paragraph spacing
        try:
            p.line_spacing = 1.15  # Adjust line spacing between bullet points
            # Add space after each paragraph
            p._pPr.spcAft = 360000  # This adds 1 line of space after each point
        except:
            pass
        
        # Apply special formatting based on format type
        if format_type == 'qa':
            # Format first bullet differently if it's a question
            if i == 0 and point_text.endswith('?'):
                p.font.bold = True
                p.font.size = Pt(18)
                continue  # Skip bullet for question
                
        elif format_type == 'tables' and 'table:' in point_text.lower():
            # Format table descriptions with special styling
            p.font.italic = True
        
        elif format_type == 'diagrams' and any(term in point_text.lower() for term in ['diagram:', 'chart:', 'graph:']):
            # Format diagram descriptions
            p.font.italic = True
            p.font.color.rgb = RGBColor(0, 112, 192)  # Blue color
            
        # Set consistent font size
        p.font.size = Pt(18)
        
        # Only add bullets for regular points (not headers/questions)
        if not point_text.endswith((':', '?')):
            try:
                # Get paragraph properties
                pPr = p._pPr
                
                # Set bullet properties using PowerPoint's built-in formatting
                pPr.get_or_add_buNone()  # Clear any existing bullet format
                pPr.get_or_add_buFont()  # Add bullet font
                pPr.get_or_add_buAutoNum()  # Add auto-numbering properties
                
                # Set bullet position and indentation
                pPr.marL = 342900  # Left margin for bullet points (0.95cm)
                pPr.indent = -342900  # Negative indent for hanging bullet points
                pPr.lvl = 0  # First level bullet
                
            except Exception as e:
                print(f"Warning: Could not set bullet properties for point: {point_text[:30]}...")
    
    return slide

def create_qa_slide(prs, question, answers):
    """Create a slide in Question & Answer format"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set question as the title
    title = slide.shapes.title
    title.text = question
    
    # Add answers as bullet points
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for answer in answers:
        answer_text = answer.strip()
        if answer_text.startswith('- '):
            answer_text = answer_text[2:]
            
        p = tf.add_paragraph()
        p.text = answer_text
        p.level = 0
    
    return slide

def create_diagram_slide(prs, slide_title, diagram_description, page_number=None):
    """Create a slide with a placeholder for a diagram"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Extract diagram type and description
    diagram_type = "Diagram"
    diagram_content = diagram_description
    
    # Try to identify diagram type
    diagram_types = ["flowchart", "pie chart", "bar chart", "line graph", "org chart", 
                    "venn diagram", "timeline", "mind map", "scatter plot", "diagram"]
    
    for d_type in diagram_types:
        if d_type in diagram_description.lower():
            diagram_type = d_type.title()
            break
    
    # Look for images that might match this diagram description
    images_dir = os.path.join('extract', 'images')
    diagram_images = []
    
    # Only look for specific page images if we have a page number
    if page_number and os.path.exists(images_dir):
        # Look specifically for images from this page, supporting both PDF extraction and topic generator
        page_images = [f for f in os.listdir(images_dir) 
                      if (f'page_{page_number}_img_' in f.lower() or f'topic_{page_number}_img_' in f.lower()) and 
                      any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
        
        # Get image titles to help find the most relevant one
        image_titles = get_image_titles()
        
        # First try: Find images with titles matching the diagram type
        for img_file in page_images:
            if 'img_' in img_file:
                match = re.search(r'img_(\d+)', img_file)
                if match:
                    img_index = match.group(1)
                    
                    # Handle both formats
                    if 'page_' in img_file:
                        img_key = f"page_{page_number}_img_{img_index}"
                    elif 'topic_' in img_file:
                        img_key = f"page_{page_number}_img_{img_index}"
                    else:
                        continue
                        
                    if img_key in image_titles:
                        title = image_titles[img_key].lower()
                        # Check if the image title contains diagram type keywords
                        if any(d_type.lower() in title for d_type in diagram_types):
                            diagram_images.append(img_file)
                        # Also check if the diagram description contains words from the image title
                        elif any(word in diagram_description.lower() for word in title.split() if len(word) > 3):
                            diagram_images.append(img_file)
        
        # If we didn't find any matching images, just use the first image from this page
        if not diagram_images and page_images:
            diagram_images = [page_images[0]]
    
    # If no page-specific images found, look for any images
    if not diagram_images and os.path.exists(images_dir):
        all_images = [f for f in os.listdir(images_dir) 
                    if any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
        if all_images:
            # Just use the first available image
            diagram_images = [all_images[0]]
    
    # If we have images, use the first one
    if diagram_images:
        try:
            # Add the image in the center
            image_path = os.path.join(images_dir, diagram_images[0])
            slide.shapes.add_picture(
                image_path,
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(4)
            )
            
            # Get caption for the image
            img_caption = ""
            if page_number:
                img_file = diagram_images[0]
                if 'img_' in img_file:
                    match = re.search(r'img_(\d+)', img_file)
                    if match:
                        img_index = match.group(1)
                        
                        # Handle both formats
                        if 'page_' in img_file:
                            img_key = f"page_{page_number}_img_{img_index}"
                        elif 'topic_' in img_file:
                            img_key = f"page_{page_number}_img_{img_index}"
                        else:
                            img_key = None
                            
                        image_titles = get_image_titles()
                        if img_key and img_key in image_titles:
                            img_caption = image_titles[img_key]
            
            # Add description below with caption if available
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            
            p = desc_frame.add_paragraph()
            # Use image caption if available, otherwise use the diagram description
            display_text = img_caption if img_caption else diagram_content
            p.text = display_text
            p.font.size = Pt(12)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            
            print(f"Added diagram slide with image: {diagram_images[0]}")
            return slide
        except Exception as e:
            print(f"Error adding diagram image: {str(e)}")
    
    # If no images or error, add placeholder
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    diagram_frame = diagram_box.text_frame
    diagram_frame.word_wrap = True
    
    p = diagram_frame.add_paragraph()
    p.text = f"[{diagram_type} Placeholder]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    
    p = desc_frame.add_paragraph()
    p.text = diagram_content
    p.font.size = Pt(12)
    p.font.italic = True
    
    return slide

def create_powerpoint(slides_data, output_dir, theme_key='theme1', creator_name="Aditya"):
    """Create PowerPoint with specified theme."""
    print(f"\nCreating PowerPoint presentation...")
    print(f"Theme: {theme_key}")
    print(f"Output directory: {output_dir}")
    
    # Get theme configuration
    theme = THEMES.get(theme_key, THEMES['theme1'])
    theme_path = os.path.join('themes', theme['file'])
    print(f"Theme path: {theme_path}")
    
    # Create presentation with selected theme
    try:
        prs = Presentation(theme_path)
        print("Successfully loaded theme")
    except Exception as e:
        print(f"Error loading theme: {str(e)}")
        print("Falling back to blank presentation")
        prs = Presentation()
    
    # Detect format type from the content
    format_type = detect_format_type(slides_data)
    print(f"Detected presentation format: {format_type}")
    
    # Create title slide (first slide)
    print("\nCreating title slide...")
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = slides_data[0]['title']
    
    # Add creator name as subtitle
    if hasattr(slide, 'placeholders'):
        subtitle = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:
                subtitle = shape
                break
        
        if subtitle:
            subtitle.text = f"Created by: {creator_name}"
        else:
            left = Inches(1)
            top = Inches(4)
            width = Inches(8)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = f"Created by: {creator_name}"
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            textbox.text_frame.paragraphs[0].font.size = Pt(18)
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
    
    # Create index slide (second slide)
    print("\nCreating index slide...")
    slide_layout = prs.slide_layouts[1]
    index_slide = prs.slides.add_slide(slide_layout)
    title = index_slide.shapes.title
    title.text = "Index"
    
    content = index_slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Add all slide titles to the index
    for i, slide_data in enumerate(slides_data[2:], start=1):
        p = tf.add_paragraph()
        p.text = f"{i}. {slide_data['title']}"
        p.font.size = Pt(18)
        p.level = 0
        try:
            p.line_spacing = 1.15
        except:
            pass
    
    # Get all images and categorize them
    images_dir = os.path.join('extract', 'images')
    pdf_images = {}  # Dictionary to store PDF images by page number
    unmapped_images = []  # List to store all images that can't be mapped
    used_images = set()  # Keep track of which images have been used
    
    if os.path.exists(images_dir):
        try:
            for img_file in os.listdir(images_dir):
                if any(img_file.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                    # Check if it's a PDF-extracted image with valid page number
                    pdf_match = re.search(r'page_(\d+)_img_(\d+)', img_file)
                    if pdf_match:
                        page_num = int(pdf_match.group(1))
                        # Only map images to valid slide numbers
                        if 3 <= page_num <= len(slides_data) + 2:  # +2 for title and index slides
                            if page_num not in pdf_images:
                                pdf_images[page_num] = []
                            pdf_images[page_num].append(img_file)
                        else:
                            unmapped_images.append(img_file)
                    else:
                        unmapped_images.append(img_file)
            
            print(f"Found {sum(len(images) for images in pdf_images.values())} mapped PDF images and {len(unmapped_images)} unmapped images")
        except Exception as e:
            print(f"Error processing images directory: {str(e)}")
            traceback.print_exc()

    # Create content slides with their mapped PDF images
    print("\nCreating content slides...")
    for slide_num, slide_data in enumerate(slides_data[2:], start=3):
        print(f"\nProcessing slide {slide_num}...")
        slide_title = slide_data['title']
        
        if isinstance(slide_data['content'], list):
            bullet_points = slide_data['content']
        else:
            bullet_points = [line.strip() for line in slide_data['content'].split('\n') if line.strip()]
        
        if not bullet_points:
            print(f"Skipping empty slide {slide_num}")
            continue
            
        print(f"Creating bullet slide: {slide_title}")
        slide = create_bullet_slide(prs, slide_title, bullet_points, format_type)
        
        # Add PDF images for this slide if they exist
        if slide_num in pdf_images and pdf_images[slide_num]:
            print(f"Adding {len(pdf_images[slide_num])} PDF images for slide {slide_num}")
            
            # Create image slide
            image_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add title
            title_box = image_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = f"{slide_title} - Visual Content"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Calculate layout based on number of images
            if len(pdf_images[slide_num]) == 1:
                img_width = Inches(8)
                img_height = Inches(4.5)
                left = Inches(1)
                top = Inches(1.5)
            else:
                img_width = Inches(4)
                img_height = Inches(3)
                images_per_row = 2
            
            # Add images
            for idx, img_file in enumerate(sorted(pdf_images[slide_num])):
                try:
                    image_path = os.path.join(images_dir, img_file)
                    print(f"Adding mapped PDF image: {image_path}")
                    
                    if len(pdf_images[slide_num]) > 1:
                        row = idx // images_per_row
                        col = idx % images_per_row
                        left = Inches(1 + col * 4.5)
                        top = Inches(1.5 + row * 3.5)
                    
                    # Add image
                    pic = image_slide.shapes.add_picture(
                        image_path,
                        left=left,
                        top=top,
                        width=img_width,
                        height=img_height
                    )
                    used_images.add(img_file)
                    
                    # Get image title from image_titles.txt if available
                    titles_file = os.path.join('extract', 'image_titles.txt')
                    caption_text = None
                    if os.path.exists(titles_file):
                        with open(titles_file, 'r', encoding='utf-8') as f:
                            for line in f:
                                if img_file.split('.')[0] in line:
                                    caption_text = line.split('|')[1].strip()
                                    break
                    
                    if not caption_text:
                        caption_text = f"Figure {idx + 1}"
                    
                    # Add caption
                    caption = image_slide.shapes.add_textbox(
                        left=left,
                        top=top + img_height + Inches(0.1),
                        width=img_width,
                        height=Inches(0.3)
                    )
                    caption_frame = caption.text_frame
                    caption_frame.text = caption_text
                    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    caption_frame.paragraphs[0].font.size = Pt(10)
                    caption_frame.paragraphs[0].font.italic = True
                    print(f"Added caption: {caption_text}")
                    
                except Exception as e:
                    print(f"Error adding mapped PDF image {img_file}: {str(e)}")
                    traceback.print_exc()
    
    # Add any unmapped images at the end
    remaining_images = unmapped_images + [img for page_images in pdf_images.values() for img in page_images if img not in used_images]
    if remaining_images:
        print(f"\nAdding {len(remaining_images)} unmapped/unused images...")
        for i in range(0, len(remaining_images), 2):
            # Create a new slide for images
            image_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add title
            title_box = image_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = f"Additional Images - Part {i//2 + 1}"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            # Get current batch of images (up to 2)
            current_images = remaining_images[i:i+2]
            
            # Calculate layout
            if len(current_images) == 1:
                img_width = Inches(8)
                img_height = Inches(4.5)
                left = Inches(1)
                top = Inches(1.5)
            else:
                img_width = Inches(4)
                img_height = Inches(3)
                images_per_row = 2
            
            # Add images
            for idx, img_file in enumerate(current_images):
                try:
                    image_path = os.path.join(images_dir, img_file)
                    print(f"Adding unmapped image: {image_path}")
                    
                    if len(current_images) > 1:
                        row = idx // images_per_row
                        col = idx % images_per_row
                        left = Inches(1 + col * 4.5)
                        top = Inches(1.5 + row * 3.5)
                    
                    # Add image
                    pic = image_slide.shapes.add_picture(
                        image_path,
                        left=left,
                        top=top,
                        width=img_width,
                        height=img_height
                    )
                    
                    # Get image title from image_titles.txt if available
                    titles_file = os.path.join('extract', 'image_titles.txt')
                    caption_text = None
                    if os.path.exists(titles_file):
                        with open(titles_file, 'r', encoding='utf-8') as f:
                            for line in f:
                                if img_file.split('.')[0] in line:
                                    caption_text = line.split('|')[1].strip()
                                    break
                    
                    if not caption_text:
                        caption_text = f"Additional Figure {i + idx + 1}"
                    
                    # Add caption
                    caption = image_slide.shapes.add_textbox(
                        left=left,
                        top=top + img_height + Inches(0.1),
                        width=img_width,
                        height=Inches(0.3)
                    )
                    caption_frame = caption.text_frame
                    caption_frame.text = caption_text
                    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    caption_frame.paragraphs[0].font.size = Pt(10)
                    caption_frame.paragraphs[0].font.italic = True
                    print(f"Added caption: {caption_text}")
                    
                except Exception as e:
                    print(f"Error adding unmapped image {img_file}: {str(e)}")
                    traceback.print_exc()
    
    # Save the presentation
    if os.path.isdir(output_dir):
        output_file = os.path.join(output_dir, 'generated_presentation.pptx')
    else:
        output_file = output_dir
        
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    
    print(f"\nSaving presentation to: {output_file}")
    try:
        prs.save(output_file)
        print(f"Successfully saved presentation with theme {theme_key}")
    except Exception as e:
        print(f"Error saving presentation: {str(e)}")
        traceback.print_exc()
        try:
            print("Attempting to save without theme...")
            blank_prs = Presentation()
            blank_prs.save(output_file)
            print("Successfully saved presentation without theme")
        except Exception as e2:
            print(f"Error saving fallback presentation: {str(e2)}")
    
    return output_file

def main():
    vba_file_path = 'create_presentation.vba'
    output_dir = 'output'
    
    # List available themes
    print("\nAvailable themes:")
    for theme in list_available_themes():
        print(f"- {theme['name']} (key: {theme['key']})")
    
    slides_data = parse_vba_file(vba_file_path)
    ppt_path = create_powerpoint(slides_data, output_dir, creator_name="Aditya")
    return ppt_path

if __name__ == "__main__":
    main()
